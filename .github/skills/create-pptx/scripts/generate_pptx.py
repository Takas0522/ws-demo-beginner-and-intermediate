#!/usr/bin/env python3
"""
PowerPoint スライド生成スクリプト

スライド定義JSONをもとに、テンプレートを使用してPowerPointファイルを生成する。

Usage:
    python generate_pptx.py <slides.json> [--template <template.pptx>]
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === カラーパレット ===
BG_DARK       = RGBColor(0x1E, 0x1E, 0x2E)
BG_SURFACE    = RGBColor(0x2A, 0x2A, 0x3C)
ACCENT        = RGBColor(0x00, 0x9B, 0xF5)
ACCENT2       = RGBColor(0x00, 0xD4, 0xAA)
ACCENT3       = RGBColor(0xAA, 0x70, 0xFF)
TEXT_PRIMARY   = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_SECONDARY = RGBColor(0xA0, 0xA0, 0xB0)
DIVIDER        = RGBColor(0x3A, 0x3A, 0x4E)
IMG_PLACEHOLDER = RGBColor(0x35, 0x35, 0x48)

FONT_NAME = "メイリオ"
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    return add_shape(slide, left, top, width, height, color)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_NAME):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # テキストフレームに内部余白を設定
    txBox.text_frame.margin_left   = Pt(2)
    txBox.text_frame.margin_right  = Pt(2)
    txBox.text_frame.margin_top    = Pt(2)
    txBox.text_frame.margin_bottom = Pt(2)

    lines = text.split('\n') if text else ['']
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = _Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = alignment

        # 段落間スペース（空行は半分）
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        space_after = int(_Pt(font_size * 0.15))
        pPr.set('spc', str(int(_Pt(font_size * 0.15))))

        # 行間を1.1倍に設定
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '110000')

    return txBox


def add_header_bar(slide, title_text):
    """共通ヘッダーバー（コンテンツスライド用）"""
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BG_SURFACE)
    add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                 title_text, font_size=32, bold=True, color=TEXT_PRIMARY)


def add_footer_line(slide, footer_text=None):
    """共通フッターライン"""
    add_accent_bar(slide, Inches(0), Inches(7.1), SLIDE_WIDTH, Inches(0.03), DIVIDER)
    if footer_text:
        add_text_box(slide, Inches(0.8), Inches(7.1), Inches(4), Inches(0.4),
                     footer_text, font_size=10, color=TEXT_SECONDARY)


# === レイアウト生成関数 ===

def create_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 左アクセントバー
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ACCENT)
    # 上部バー
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    # 装飾ブロック右下
    add_shape(slide, Inches(10.5), Inches(5.5), Inches(2.8), Inches(2), BG_SURFACE)
    add_accent_bar(slide, Inches(10.5), Inches(5.5), Inches(2.8), Inches(0.04), ACCENT2)
    # タイトル
    add_text_box(slide, Inches(1.0), Inches(2.2), Inches(9), Inches(1.5),
                 data.get("title", "タイトル"), font_size=44, bold=True)
    # サブタイトル
    add_text_box(slide, Inches(1.0), Inches(3.8), Inches(9), Inches(0.8),
                 data.get("subtitle", ""), font_size=22, color=TEXT_SECONDARY)
    # 区切り線
    add_accent_bar(slide, Inches(1.0), Inches(4.8), Inches(3), Inches(0.04), ACCENT)
    # 発表者
    add_text_box(slide, Inches(1.0), Inches(5.2), Inches(6), Inches(0.6),
                 data.get("author", ""), font_size=16, color=TEXT_SECONDARY)
    return slide


def create_section_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_shape(slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(2.8), BG_SURFACE)
    add_accent_bar(slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.05), ACCENT)
    add_text_box(slide, Inches(1.0), Inches(2.7), Inches(2), Inches(1),
                 data.get("number", "01"), font_size=60, bold=True, color=ACCENT)
    add_text_box(slide, Inches(3.5), Inches(3.0), Inches(8), Inches(1),
                 data.get("title", ""), font_size=36, bold=True)
    add_text_box(slide, Inches(3.5), Inches(4.0), Inches(8), Inches(0.8),
                 data.get("description", ""), font_size=18, color=TEXT_SECONDARY)
    return slide


def create_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_header_bar(slide, data.get("title", ""))
    add_text_box(slide, Inches(0.8), Inches(1.55), Inches(11.7), Inches(5.3),
                 data.get("body", ""), font_size=18)
    add_footer_line(slide, data.get("footer"))
    return slide


def create_two_column_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_header_bar(slide, data.get("title", ""))

    # 左カラム
    add_shape(slide, Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.3), BG_SURFACE)
    add_accent_bar(slide, Inches(0.5), Inches(1.55), Inches(6.1), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(0.85), Inches(1.75), Inches(5.4), Inches(0.65),
                 data.get("left_heading", ""), font_size=20, bold=True)
    add_text_box(slide, Inches(0.85), Inches(2.55), Inches(5.4), Inches(4.0),
                 data.get("left_body", ""), font_size=15, color=TEXT_SECONDARY)

    # 右カラム
    add_shape(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(5.3), BG_SURFACE)
    add_accent_bar(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(0.05), ACCENT2)
    add_text_box(slide, Inches(7.25), Inches(1.75), Inches(5.4), Inches(0.65),
                 data.get("right_heading", ""), font_size=20, bold=True)
    add_text_box(slide, Inches(7.25), Inches(2.55), Inches(5.4), Inches(4.0),
                 data.get("right_body", ""), font_size=15, color=TEXT_SECONDARY)

    add_footer_line(slide)
    return slide


def create_three_cards_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_header_bar(slide, data.get("title", ""))

    cards = data.get("cards", [])
    accent_colors = [ACCENT, ACCENT2, ACCENT3]

    for i in range(min(len(cards), 3)):
        card = cards[i]
        color = accent_colors[i % 3]
        x = Inches(0.5 + i * 4.28)

        add_shape(slide, x, Inches(1.55), Inches(3.9), Inches(5.3), BG_SURFACE)
        add_accent_bar(slide, x, Inches(1.55), Inches(3.9), Inches(0.06), color)
        add_text_box(slide, x + Inches(0.3), Inches(1.85), Inches(1.5), Inches(0.9),
                     card.get("number", f"{i+1:02d}"), font_size=44, bold=True, color=color)
        add_text_box(slide, x + Inches(0.3), Inches(2.85), Inches(3.3), Inches(0.65),
                     card.get("heading", ""), font_size=20, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(3.65), Inches(3.3), Inches(3.0),
                     card.get("body", ""), font_size=15, color=TEXT_SECONDARY)

    add_footer_line(slide)
    return slide


def create_visual_text_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_header_bar(slide, data.get("title", ""))

    # 画像エリア
    image_path = data.get("image_path")
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.6), Inches(1.8), Inches(6), Inches(5))
    else:
        add_shape(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(5), IMG_PLACEHOLDER)
        add_text_box(slide, Inches(2.0), Inches(3.8), Inches(3.5), Inches(1),
                     data.get("image_placeholder", "画像・図表エリア"),
                     font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # テキストエリア
    add_text_box(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.8),
                 data.get("content_heading", ""), font_size=24, bold=True)
    add_accent_bar(slide, Inches(7.2), Inches(2.9), Inches(2), Inches(0.04), ACCENT)
    add_text_box(slide, Inches(7.2), Inches(3.3), Inches(5.5), Inches(3.5),
                 data.get("content_body", ""), font_size=16, color=TEXT_SECONDARY)

    add_footer_line(slide)
    return slide


def create_ending_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.2),
                 data.get("message", "Thank You"), font_size=54, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.04), ACCENT)
    add_text_box(slide, Inches(0), Inches(4.3), SLIDE_WIDTH, Inches(0.8),
                 data.get("submessage", "ご清聴ありがとうございました"),
                 font_size=22, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.5),
                 data.get("contact", ""), font_size=14, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(0), Inches(6.8), Inches(4), Inches(0.7), BG_SURFACE)
    add_accent_bar(slide, Inches(0), Inches(6.8), Inches(4), Inches(0.04), ACCENT2)
    return slide


def create_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_header_bar(slide, data.get("title", ""))

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    col_count = len(headers) if headers else (len(rows[0]) if rows else 1)
    row_count = len(rows) + 1  # +1 for header

    table_left   = Inches(0.5)
    table_top    = Inches(1.55)
    table_width  = Inches(12.3)
    row_h        = Inches(max(0.42, 4.8 / row_count))
    table_height = row_h * row_count

    table_shape = slide.shapes.add_table(row_count, col_count,
                                         table_left, table_top,
                                         table_width, int(table_height))
    table = table_shape.table

    # 行高を均等に設定
    for ri in range(row_count):
        table.rows[ri].height = int(row_h)

    # ヘッダー行スタイル
    for ci, header_text in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(header_text)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_PRIMARY
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x6A, 0xAA)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # データ行スタイル
    for ri, row_data in enumerate(rows):
        bg = BG_SURFACE if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x34)
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_NAME
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = TEXT_PRIMARY
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # テーブル罫線
    from pptx.oxml.ns import qn
    from lxml import etree
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tc_pr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tc_pr.find(qn(border_name))
                if ln is None:
                    ln = etree.SubElement(tc_pr, qn(border_name), w='9525', cap='flat', cmpd='sng')
                else:
                    ln.set('w', '9525')
                sf = ln.find(qn('a:solidFill'))
                if sf is None:
                    sf = etree.SubElement(ln, qn('a:solidFill'))
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is None:
                    srgb = etree.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', '3A3A4E')

    add_footer_line(slide, data.get("footer"))
    return slide


LAYOUT_MAP = {
    "title": create_title_slide,
    "section": create_section_slide,
    "content": create_content_slide,
    "two-column": create_two_column_slide,
    "three-cards": create_three_cards_slide,
    "visual-text": create_visual_text_slide,
    "table": create_table_slide,
    "ending": create_ending_slide,
}


def main():
    parser = argparse.ArgumentParser(description="スライド定義JSONからPowerPointを生成")
    parser.add_argument("slides_json", help="スライド定義JSONファイルのパス")
    parser.add_argument("--template", default=None,
                        help="テンプレートPPTXのパス（未指定時はデフォルトテンプレートを使用）")
    args = parser.parse_args()

    # JSON読み込み
    with open(args.slides_json, "r", encoding="utf-8") as f:
        config = json.load(f)

    # プレゼンテーション作成（テンプレートは装飾の参考だが、Blankレイアウトのみ使う）
    template_path = args.template
    if template_path is None:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(script_dir, "..", "resources", "template_dark.pptx")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # テンプレートの既存スライドを削除（レイアウト定義のみ残す）
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    # スライド生成
    slides_data = config.get("slides", [])
    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        creator = LAYOUT_MAP.get(layout)
        if creator is None:
            print(f"警告: 不明なレイアウト '{layout}' をスキップしました", file=sys.stderr)
            continue
        creator(prs, slide_data)

    # 出力
    output_path = config.get("output", "output.pptx")
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"生成完了: {output_path} ({len(slides_data)} スライド)")


if __name__ == "__main__":
    main()
